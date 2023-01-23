# imports 
import seaborn as sns
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.lines as mlines
from matplotlib import rcParams
import pandas as pd
import collections.abc
from pptx import Presentation
import urllib.request
from collections import Counter
from datetime import datetime
import ETL_path_manager as ETL_pm

# global variables
PRS = Presentation(ETL_pm.get_presentation_template())
TITLE_SLIDE_LAYOUT = PRS.slide_masters[-1]

def write_shape_text(shape, text):
    """
    writes text within a shape object
    every textbox, picture, wordart, etc. is a type of shape object
    """
    
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text

def delete_shape(slide, idx):
    slide.shapes.element.remove(slide.shapes[idx].element)

def insert_picture(shape, path):
    shape.insert_picture(path)

def modify_title_slide(playlist_name):
    """
    program refernces a template ppt when creating a ppt
    that template contains a title slide
    this function modifies said title slide
    title slide = name of playlist, picture of playlist and date/time
    """
    
    slide = PRS.slides[0]

    text = playlist_name + ' @ ' + datetime.now().strftime("%H:%M:%S %d/%m/%Y")
    write_shape_text(slide.shapes[0], text)
    
    write_shape_text(slide.shapes[1], 'By Timothy Cameron')
    
    name = 'playlist_image.png'
    path = ETL_pm.get_path(playlist_name, name, 'images')
    insert_picture(slide.shapes[2], path)

def create_section_header(title, subtitle=None):
    slide = PRS.slides.add_slide(TITLE_SLIDE_LAYOUT.slide_layouts[2])
    
    write_shape_text(slide.shapes[0], title)
    
    if subtitle == None:
        delete_shape(slide, 1)
    else:
        write_shape_text(slide.shapes[1], subtitle)
 
def define_y_axis(key):
    """
    defines the y-axis for the box plot
    different audio features possess different scales
    hence the box plot has to be adapated to fit
    """
    
    zero_to_one_keys = ['danceability', 
                        'energy', 
                        'speechiness', 
                        'acousticness', 
                        'liveness',
                        'valence']
    
    if key in zero_to_one_keys:
        y_lower = 0
        y_upper = 1
    elif key == 'instrumentalness':
        y_lower = None
        y_upper = None
    elif key == 'loudness':
        y_lower = -60
        y_upper = 0
        
    return y_lower, y_upper

def create_boxplot(audio_feature, key, playlist_name):
    """
    creates a box plot for the audio features and saves them in a subfolder
    returns a path for the ppt to reference so that the boxplot can be inserted
    """
    
    y_lower, y_upper = define_y_axis(key)
    
    plt.figure(figsize=(6.752,5.909))
    
    fig = sns.boxplot(data=audio_feature,
                      orient='v',
                      color='#1DB954', # spotify green 
                      width=0.5,
                      showmeans=True,
                      meanline=True,
                      meanprops={'color': 'black', 'linestyle': '-.'})
    
    fig.set(ylim=(y_lower,y_upper), xlim=(-0.5,0.5), xticklabels=[])
    
    mean_line_label = 'Mean = ' + str("%.2f" % np.mean(audio_feature))
    median_line_label = 'Median = ' + str("%.2f" % np.median(audio_feature))
    
    mean_line= mlines.Line2D([], [], label=mean_line_label, linestyle='-.' )
    median_line = mlines.Line2D([], [], label=median_line_label, linestyle='-') 
    fig.legend(handles=[mean_line, median_line], loc='upper right')
    fig.legend(handles=[mean_line, median_line])
    
    name = key + '_boxplot.png'
    path = ETL_pm.get_path(playlist_name, name, 'plots')
    
    plt.savefig(path)
    plt.clf()
    
    return path

def get_album_pic(playlist_name, idx, df):
    
    name = 'song' + str(idx + 1) + '.png'
    path = ETL_pm.get_path(playlist_name, name, 'images')
    
    urllib.request.urlretrieve(df.iloc[idx]['album_image_URL'], path)
    
    return path

def create_top_five_slides(df, playlist_name):
    """creates a slide for each of the top 5 songs for a given playlist
    displays: song image, song name, sort artists and album name
    the top 5 songs are the ones with the most listens for a given day
    """
    
    for idx in range(5):
        slide = PRS.slides.add_slide(TITLE_SLIDE_LAYOUT.slide_layouts[8])
        
        path = get_album_pic(playlist_name, idx, df)
        
        text = str(idx + 1) + ' - ' + df.iloc[idx]['track_name']
        write_shape_text(slide.shapes[0], text)
        
        text = 'By ' + ', '.join([artist for artist in df.iloc[idx]['artists']])
        text = text + '. Album - ' + df.iloc[idx]['album_name']
        write_shape_text(slide.shapes[2], text)
        
        insert_picture(slide.shapes[1], path)

def get_top_five_genres(df):
    """
    finds out what the most popular genres are for a given playlist
    please note that spotify doesn't assign genres to songs
    but to artists
    """
    
    top_five_sorted_genres = []
    
    # Counter method only works with flat lists
    flat_genres = [item for sublist in df['artist_genres'] for item in sublist]
    
    top_five_genres = Counter(flat_genres).most_common(5)
    
    for genre, occurence in top_five_genres:
        for _ in range(occurence):
            top_five_sorted_genres.append(genre.capitalize())
    
    return top_five_sorted_genres

def create_genre_countplot(df, playlist_name):
    """creates a countplot for the top 5 popular genres in a given playlist
    returns a path for the ppt to reference so that the countplot can be inserted
    """
    
    top_five_sorted_genres = get_top_five_genres(df)
    
    plt.figure(figsize=(11.492,4.381))
    sns.countplot(x=top_five_sorted_genres)
    
    name = 'genre_countplot.png'
    path = ETL_pm.get_path(playlist_name, name, 'plots')
    
    plt.savefig(path)
    plt.clf() 
    
    return path

def create_genre_slide(df, playlist_name):
    """creates a slide with a countplot, counting the top 5 most popular genres for a given playlist"""
    
    slide = PRS.slides.add_slide(TITLE_SLIDE_LAYOUT.slide_layouts[1])
    
    write_shape_text(slide.shapes[0], 'Most Popular Genres')
    
    path = create_genre_countplot(df, playlist_name)
    insert_picture(slide.shapes[1], path)
    
def create_audio_feature_slides(df, playlist_name):
    """creates a slide for each spotify audio feature
    each slide contains the name of the audio feature with its associated description
    and a boxplot of the data
    """
    
    audio_features = {
        'danceability': """Danceability describes how \
suitable a track is for \
dancing based on a combination of musical elements \
including tempo, rhythm stability, beat strength, \
and overall regularity. A value of 0.0 is least \
danceable and 1.0 is most danceable.\
        """,
        
        'energy': """Energy is a measure from 0.0 to 1.0 and \
represents a perceptual measure of intensity and activity. Typically, \
energetic tracks feel fast, loud, and noisy. For example, \
death metal has high energy, while a Bach prelude scores \
low on the scale. Perceptual features contributing to this \
attribute include dynamic range, perceived loudness, timbre, \
onset rate, and general entropy. \
        """,
        
        'loudness': """The overall loudness of a track in decibels (dB). \
Loudness values are averaged across the entire track \
and are useful for comparing relative loudness of tracks. \
Loudness is the quality of a sound that is the primary \
psychological correlate of physical strength (amplitude). \
Values typically range between -60 and 0 db. \
        """,
        
        'speechiness': """Speechiness detects the presence of spoken words in a track. \
The more exclusively speech-like the recording \
(e.g. talk show, audio book, poetry), \
the closer to 1.0 the attribute value. Values above 0.66 \
describe tracks that are probably made entirely of spoken words. \
Values between 0.33 and 0.66 describe tracks that may contain \
both music and speech, either in sections or layered, \
including such cases as rap music. Values below 0.33 most \
likely represent music and other non-speech-like tracks. \
        """,
        
        'acousticness': """A confidence measure from 0.0 to 1.0 of whether the track \
is acoustic. 1.0 represents high confidence the track is acoustic. \
        """,
        
        'instrumentalness': """Predicts whether a track contains no vocals. "Ooh" and "aah" \
sounds are treated as instrumental in this context. Rap \
or spoken word tracks are clearly "vocal". The closer the \
instrumentalness value is to 1.0, the greater likelihood \
the track contains no vocal content. Values above 0.5 are \
intended to represent instrumental tracks, but confidence \
is higher as the value approaches 1.0. \
        """,
        
        'liveness': """Detects the presence of an audience in the recording. \
Higher liveness values represent an increased probability \
that the track was performed live. A value above 0.8 \
provides strong likelihood that the track is live. \
        """,
        
        'valence': """A measure from 0.0 to 1.0 describing the musical positiveness \
conveyed by a track. Tracks with high valence sound more \
positive (e.g. happy, cheerful, euphoric), while tracks \
with low valence sound more negative (e.g. sad, depressed, angry). \
        """
    }

    for key in audio_features.keys():
        slide = PRS.slides.add_slide(TITLE_SLIDE_LAYOUT.slide_layouts[8])
        
        write_shape_text(slide.shapes[0], key.capitalize())
        
        write_shape_text(slide.shapes[2], audio_features[key])
        
        path = create_boxplot(df[key], key, playlist_name)
        insert_picture(slide.shapes[1], path)

def create_powerpoint(df, playlist_name):
    """creates a powerpoint giving a summary of the information for a given daily top 50 playlist
    the slides are created/edited sequentially with the below function calls
    the slide is then saved to a sub folder
    """

    modify_title_slide(playlist_name)
    
    create_section_header('Top 5 Songs')

    create_top_five_slides(df, playlist_name)

    create_genre_slide(df, playlist_name)

    create_section_header('Audio Features')

    create_audio_feature_slides(df, playlist_name)

    text = 'Please refer to ' + playlist_name + '.csv for the extracted results.'
    create_section_header('Thank you!', text)

    name = playlist_name + '.pptx'
    path = ETL_pm.get_path(playlist_name, name)
    PRS.save(path)