import pandas 
from pptx import Presentation
from datetime import datetime
from ETL_path_manager_scripts.ETL_path_manager import get_presentation_template
from ETL_path_manager_scripts.ETL_path_manager import get_path
from powerpoint_creator_scripts.write_shape_text import write_shape_text
from powerpoint_creator_scripts.insert_picture import insert_picture
from powerpoint_creator_scripts.delete_shape import delete_shape
from powerpoint_creator_scripts.get_album_pic import get_album_pic
from powerpoint_creator_scripts.create_genre_countplot import create_genre_countplot
from powerpoint_creator_scripts.create_boxplot import create_boxplot

PRS = Presentation(get_presentation_template())
TITLE_SLIDE_LAYOUT = PRS.slide_masters[-1]

def modify_title_slide(playlist_name):
    """
    set name of playlist, picture of playlist and current date/time on a title slide
    """
    
    slide = PRS.slides[0]

    text = playlist_name + ' @ ' + datetime.now().strftime("%H:%M:%S %d/%m/%Y")
    write_shape_text(slide.shapes[0], text)
    
    write_shape_text(slide.shapes[1], 'By Timothy Cameron')
    
    path = get_path(playlist_name, 'playlist_image.png', 'images')
    insert_picture(slide.shapes[2], path)

def create_section_header(title, subtitle=None):
    slide = PRS.slides.add_slide(TITLE_SLIDE_LAYOUT.slide_layouts[2])
    
    write_shape_text(slide.shapes[0], title)
    
    if subtitle == None:
        delete_shape(slide, 1)
    else:
        write_shape_text(slide.shapes[1], subtitle)

def create_top_five_slides(df, playlist_name):
    """creates a slide for each of the top 5 songs for a given playlist
    displays: song image, song name, sort artists and album name
    the top 5 songs are the ones with the most listens for a given day
    """
    
    for idx in range(5):
        slide = PRS.slides.add_slide(TITLE_SLIDE_LAYOUT.slide_layouts[8])
        
        path = get_album_pic(playlist_name, idx, df)
        insert_picture(slide.shapes[1], path)
        
        text = str(idx + 1) + ' - ' + df.iloc[idx]['track_name']
        write_shape_text(slide.shapes[0], text)
        
        text = 'By ' + ', '.join([artist for artist in df.iloc[idx]['artists']])
        text = text + '. Album - ' + df.iloc[idx]['album_name']
        write_shape_text(slide.shapes[2], text)

def create_genre_slide(df, playlist_name):
    """creates a slide with a countplot, counting the top 5 most popular genres for a given playlist"""
    
    slide = PRS.slides.add_slide(TITLE_SLIDE_LAYOUT.slide_layouts[1])
    
    write_shape_text(slide.shapes[0], 'Most Popular Genres')
    
    path = create_genre_countplot(df, playlist_name)
    insert_picture(slide.shapes[1], path)
    
def create_audio_feature_slides(df, playlist_name):
    """creates a slide for each spotify audio feature
    each slide contains the name of the audio feature with its associated description
    and a boxplot of the data
    """
    
    # make this a separate imported JSON file
    audio_features = {
        'danceability': """Danceability describes how \
suitable a track is for \
dancing based on a combination of musical elements \
including tempo, rhythm stability, beat strength, \
and overall regularity. A value of 0.0 is least \
danceable and 1.0 is most danceable.\
        """,
        
        'energy': """Energy is a measure from 0.0 to 1.0 and \
represents a perceptual measure of intensity and activity. Typically, \
energetic tracks feel fast, loud, and noisy. For example, \
death metal has high energy, while a Bach prelude scores \
low on the scale. Perceptual features contributing to this \
attribute include dynamic range, perceived loudness, timbre, \
onset rate, and general entropy. \
        """,
        
        'loudness': """The overall loudness of a track in decibels (dB). \
Loudness values are averaged across the entire track \
and are useful for comparing relative loudness of tracks. \
Loudness is the quality of a sound that is the primary \
psychological correlate of physical strength (amplitude). \
Values typically range between -60 and 0 db. \
        """,
        
        'speechiness': """Speechiness detects the presence of spoken words in a track. \
The more exclusively speech-like the recording \
(e.g. talk show, audio book, poetry), \
the closer to 1.0 the attribute value. Values above 0.66 \
describe tracks that are probably made entirely of spoken words. \
Values between 0.33 and 0.66 describe tracks that may contain \
both music and speech, either in sections or layered, \
including such cases as rap music. Values below 0.33 most \
likely represent music and other non-speech-like tracks. \
        """,
        
        'acousticness': """A confidence measure from 0.0 to 1.0 of whether the track \
is acoustic. 1.0 represents high confidence the track is acoustic. \
        """,
        
        'instrumentalness': """Predicts whether a track contains no vocals. "Ooh" and "aah" \
sounds are treated as instrumental in this context. Rap \
or spoken word tracks are clearly "vocal". The closer the \
instrumentalness value is to 1.0, the greater likelihood \
the track contains no vocal content. Values above 0.5 are \
intended to represent instrumental tracks, but confidence \
is higher as the value approaches 1.0. \
        """,
        
        'liveness': """Detects the presence of an audience in the recording. \
Higher liveness values represent an increased probability \
that the track was performed live. A value above 0.8 \
provides strong likelihood that the track is live. \
        """,
        
        'valence': """A measure from 0.0 to 1.0 describing the musical positiveness \
conveyed by a track. Tracks with high valence sound more \
positive (e.g. happy, cheerful, euphoric), while tracks \
with low valence sound more negative (e.g. sad, depressed, angry). \
        """
    }

    for key in audio_features:
        slide = PRS.slides.add_slide(TITLE_SLIDE_LAYOUT.slide_layouts[8])
        
        write_shape_text(slide.shapes[0], key.capitalize())
        
        write_shape_text(slide.shapes[2], audio_features[key])
        
        path = create_boxplot(df[key], key, playlist_name)
        insert_picture(slide.shapes[1], path)

def create_powerpoint(df, playlist_name):
    """creates a powerpoint giving a summary of the information for a given daily top 50 playlist
    the slides are created/edited sequentially with the below function calls
    the slide is then saved to a sub folder
    """

    modify_title_slide(playlist_name)
    
    create_section_header('Top 5 Songs')

    create_top_five_slides(df, playlist_name)

    create_genre_slide(df, playlist_name)

    create_section_header('Audio Features')

    create_audio_feature_slides(df, playlist_name)

    text = 'Please refer to ' + playlist_name + '.csv for the extracted results.'
    create_section_header('Thank you!', text)

    name = playlist_name + '.pptx'
    path = get_path(playlist_name, name)
    PRS.save(path)
