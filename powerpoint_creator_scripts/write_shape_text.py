from pptx import Presentation

def write_shape_text(shape, text):
    """
    writes text within a shape object
    every textbox, picture, wordart, etc. is a type of shape object
    """
    
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
